import os
import unicodedata
import zipfile
import re
from typing import Optional
from typing import Optional, List
from pathlib import Path

# PDF
import fitz  # PyMuPDF - AGPL lizenziert

# Word
from docx import Document as DocxDocument

# Excel
import openpyxl
import csv


# PowerPoint
from pptx import Presentation

from project_data import ProjectData


class ReaderService:
    """
    Service zum Lesen von Text aus verschiedenen Dateiformaten.
    Nutzt spezialisierte Bibliotheken für jedes Format.

    Verwendete Bibliotheken:
    - PyMuPDF (AGPL v3.0): https://pymupdf.readthedocs.io/
    - python-docx (MIT): Für Word-Dokumente
    - openpyxl (MIT): Für Excel-Dateien
    - python-pptx (MIT): Für PowerPoint-Präsentationen
    """

    def __init__(self):
        # Unterstützte Dateiformate
        self.supported_extensions = {
            # Dokumente
            '.pdf', '.txt', '.md', '.html', '.htm',
            '.docx', '.doc', '.odt', '.rtf',
            # Tabellen
            '.xlsx', '.xls', '.csv',
            # Präsentationen
            '.pptx', '.ppt',
            # Apple Office
            '.pages', '.numbers'
        }

        # Für Statistiken/Logging
        self.stats = {
            'success': 0,
            'failed': 0,
            'unsupported': 0
        }


    def extract_text(self, filepath: str, max_chars: Optional[int] = ProjectData.search_depth) -> Optional[str]:
        """
        Zentrale Methode: erkennt Dateiformat und ruft passende Extraktionsmethode auf.

        Args:
            filepath: Pfad zur Datei
            max_chars: Maximale Zeichenanzahl - WICHTIG: Dies ist die SUCHTIEFE!
                      Sobald diese Zeichenanzahl erreicht ist, wird abgebrochen.
                      außer bei pdf da liest er immer die 1. Seite

        Returns:
            Extrahierter Text oder None bei Fehler
        """
        # Pfad normalisieren (nur für Dateisystem-Zugriff)
        path = Path(filepath)
        safe_path = path.as_posix()

        if not os.path.exists(safe_path):
            print(f"❓ Datei nicht gefunden: {safe_path}")
            return None

        ext = os.path.splitext(filepath)[1].lower()
        if ext not in self.supported_extensions:
            print(f"⚠️ Format nicht unterstützt: {filepath}")
            self.stats['unsupported'] += 1
            return None

        # Format-spezifische Extraktion
        try:
            if ext == '.pdf':
                text = self._extract_pdf(safe_path, max_chars)
            elif ext in {'.txt', '.md'}:
                text = self._extract_text_file(safe_path, max_chars)
            elif ext in {'.html', '.htm'}:
                text = self._extract_html(safe_path, max_chars)
            elif ext in {'.docx', '.doc'}:
                text = self._extract_word(safe_path, max_chars)
            elif ext in {'.xlsx', '.xls'}:
                text = self._extract_excel(safe_path, max_chars)
            elif ext == '.csv':
                text = self._extract_csv(safe_path, max_chars)
            elif ext in {'.pptx', '.ppt'}:
                text = self._extract_powerpoint(safe_path, max_chars)
            elif ext == '.odt':
                text = self._extract_opendocument(safe_path, max_chars)
            elif ext == '.rtf':
                text = self._extract_rtf(safe_path, max_chars)
            elif ext == '.numbers' or ext == '.pages':
                text = self._extract_iwork_file(safe_path, max_chars)
            else:
                text = self._extract_generic(safe_path, max_chars)

            if not text:
                return None

            # 🔥 Text für Vergleich normalisieren (NFC)
            text = unicodedata.normalize('NFC', text)

            # Text bereinigen (mehrfache Leerzeichen entfernen)
            text = ' '.join(text.split())

            self.stats['success'] += 1
            return text

        except Exception as e:
            print(f"❌ Fehler beim Lesen von {filepath}: {e}")
            self.stats['failed'] += 1
            return None

    # ===== PDF =====

    def _extract_pdf(self, filepath: str, max_chars: Optional[int] = None) -> str:
        """PDF-Extraktion mit PyMuPDF - erste Seite immer komplett, danach max_chars beachten."""
        if not os.path.exists(filepath):
            raise FileNotFoundError(f"PDF-Datei nicht gefunden: {filepath}")

        doc = None
        try:
            doc = fitz.open(filepath)
            text_parts = []
            total_chars = 0

            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                page_text = page.get_text()

                # Prüfen, ob überhaupt noch Platz ist (für Performance)
                if max_chars is not None and total_chars >= max_chars:
                    break

                if page_num == 0:
                    # erste Seite immer komplett lesen
                    text_parts.append(page_text)
                    total_chars += len(page_text)
                else:
                    # Für Folgeseiten: max_chars berücksichtigen
                    if max_chars is not None:
                        remaining = max_chars - total_chars
                        if remaining <= 0:
                            break

                        if len(page_text) > remaining:
                            # Nur Teil der Seite lesen
                            page_text = page_text[:remaining]
                            text_parts.append(page_text)
                            total_chars += remaining
                            break  # Abbruch, da max_chars erreicht
                        else:
                            # Ganze Seite passt noch
                            text_parts.append(page_text)
                            total_chars += len(page_text)
                    else:
                        # Kein Limit - ganze Seite lesen
                        text_parts.append(page_text)
                        total_chars += len(page_text)

            return "\n".join(text_parts)

        except Exception as e:
            print(f"⚠️ PyMuPDF Fehler beim Lesen von {filepath}: {e}")
            raise
        finally:
            if doc:
                doc.close()



    # ===== Textdateien =====
    def _extract_text_file(self, filepath: str, max_chars: Optional[int] = None) -> str:
        """Einfache Textdateien - bricht bei max_chars ab."""
        try:
            if max_chars:
                # Nur die ersten max_chars Zeichen lesen
                with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                    return f.read(max_chars)
            else:
                with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                    return f.read()
        except Exception as e:
            print(f"⚠️ Textdatei-Fehler: {e}")
            raise

    def _extract_html(self, filepath: str, max_chars: Optional[int] = None) -> str:
        """HTML-Dateien - bricht bei max_chars ab."""
        try:
            # Erst gesamten Text extrahieren, dann kürzen
            from bs4 import BeautifulSoup
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                soup = BeautifulSoup(f.read(), 'html.parser')
                text = soup.get_text(separator=' ', strip=True)

            if max_chars and len(text) > max_chars:
                text = text[:max_chars]
            return text

        except ImportError:
            # Fallback: Einfaches Regex-Stripping
            import re
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                html = f.read()
                text = re.sub(r'<[^>]+>', ' ', html)
                text = re.sub(r'\s+', ' ', text)

            if max_chars and len(text) > max_chars:
                text = text[:max_chars]
            return text.strip()
        except Exception as e:
            print(f"⚠️ HTML-Fehler: {e}")
            raise

    # ===== Word =====
    def _extract_word(self, filepath: str, max_chars: Optional[int] = None) -> str:
        """Word-Dokumente (.docx) - bricht bei max_chars ab."""
        try:
            doc = DocxDocument(filepath)
            text_parts = []
            total_chars = 0

            # Absätze
            for para in doc.paragraphs:
                if max_chars and total_chars >= max_chars:
                    break

                if para.text:
                    if max_chars:
                        remaining = max_chars - total_chars
                        if len(para.text) > remaining:
                            text_parts.append(para.text[:remaining])
                            total_chars = max_chars
                        else:
                            text_parts.append(para.text)
                            total_chars += len(para.text)
                    else:
                        text_parts.append(para.text)
                        total_chars += len(para.text)

            # Tabellen (nur wenn noch Platz)
            if not max_chars or total_chars < max_chars:
                for table in doc.tables:
                    if max_chars and total_chars >= max_chars:
                        break

                    for row in table.rows:
                        if max_chars and total_chars >= max_chars:
                            break

                        for cell in row.cells:
                            if max_chars and total_chars >= max_chars:
                                break

                            if cell.text:
                                if max_chars:
                                    remaining = max_chars - total_chars
                                    if len(cell.text) > remaining:
                                        text_parts.append(cell.text[:remaining])
                                        total_chars = max_chars
                                    else:
                                        text_parts.append(cell.text)
                                        total_chars += len(cell.text)
                                else:
                                    text_parts.append(cell.text)
                                    total_chars += len(cell.text)

            return "\n".join(text_parts)

        except ImportError:
            print("⚠️ python-docx nicht installiert. Bitte installieren: pip install python-docx")
            raise
        except Exception as e:
            print(f"⚠️ Word-Fehler: {e}")
            raise

    # ===== Excel =====
    def _extract_excel(self, filepath: str, max_chars: Optional[int] = None) -> str:
        """Excel-Dateien - bricht bei max_chars ab."""
        text_parts = []
        total_chars = 0

        try:
            if filepath.endswith('.xlsx'):
                import openpyxl
                wb = openpyxl.load_workbook(filepath, data_only=True)

                for sheet_name in wb.sheetnames:
                    if max_chars and total_chars >= max_chars:
                        break

                    sheet = wb[sheet_name]
                    sheet_header = f"--- Sheet: {sheet_name} ---"

                    if max_chars:
                        remaining = max_chars - total_chars
                        if len(sheet_header) > remaining:
                            text_parts.append(sheet_header[:remaining])
                            total_chars = max_chars
                        else:
                            text_parts.append(sheet_header)
                            total_chars += len(sheet_header)
                    else:
                        text_parts.append(sheet_header)
                        total_chars += len(sheet_header)

                    for row in sheet.iter_rows(values_only=True):
                        if max_chars and total_chars >= max_chars:
                            break

                        row_text = ' '.join([str(cell) for cell in row if cell is not None])
                        if row_text.strip():
                            if max_chars:
                                remaining = max_chars - total_chars
                                if len(row_text) > remaining:
                                    text_parts.append(row_text[:remaining])
                                    total_chars = max_chars
                                else:
                                    text_parts.append(row_text)
                                    total_chars += len(row_text)
                            else:
                                text_parts.append(row_text)
                                total_chars += len(row_text)

            elif filepath.endswith('.xls'):
                import xlrd
                wb = xlrd.open_workbook(filepath)

                for sheet_idx in range(wb.nsheets):
                    if max_chars and total_chars >= max_chars:
                        break

                    sheet = wb.sheet_by_index(sheet_idx)
                    sheet_header = f"--- Sheet: {sheet.name} ---"

                    if max_chars:
                        remaining = max_chars - total_chars
                        if len(sheet_header) > remaining:
                            text_parts.append(sheet_header[:remaining])
                            total_chars = max_chars
                        else:
                            text_parts.append(sheet_header)
                            total_chars += len(sheet_header)
                    else:
                        text_parts.append(sheet_header)
                        total_chars += len(sheet_header)

                    for row_idx in range(sheet.nrows):
                        if max_chars and total_chars >= max_chars:
                            break

                        row = sheet.row_values(row_idx)
                        row_text = ' '.join([str(cell) for cell in row if cell != ''])
                        if row_text.strip():
                            if max_chars:
                                remaining = max_chars - total_chars
                                if len(row_text) > remaining:
                                    text_parts.append(row_text[:remaining])
                                    total_chars = max_chars
                                else:
                                    text_parts.append(row_text)
                                    total_chars += len(row_text)
                            else:
                                text_parts.append(row_text)
                                total_chars += len(row_text)

            return "\n".join(text_parts)

        except ImportError as e:
            print(f"⚠️ Excel-Bibliothek fehlt: {e}. Bitte installieren: pip install openpyxl xlrd")
            raise
        except Exception as e:
            print(f"⚠️ Excel-Fehler: {e}")
            raise

    def _extract_csv(self, filepath: str, max_chars: Optional[int] = None) -> str:
        """CSV-Dateien - bricht bei max_chars ab."""
        try:
            text_parts = []
            total_chars = 0
            encodings = ['utf-8', 'latin-1', 'cp1252']

            for encoding in encodings:
                try:
                    with open(filepath, 'r', encoding=encoding) as f:
                        reader = csv.reader(f)
                        for row in reader:
                            if max_chars and total_chars >= max_chars:
                                break

                            row_text = ' '.join(row)
                            if row_text.strip():
                                if max_chars:
                                    remaining = max_chars - total_chars
                                    if len(row_text) > remaining:
                                        text_parts.append(row_text[:remaining])
                                        total_chars = max_chars
                                    else:
                                        text_parts.append(row_text)
                                        total_chars += len(row_text)
                                else:
                                    text_parts.append(row_text)
                                    total_chars += len(row_text)
                    break  # Erfolgreich gelesen
                except UnicodeDecodeError:
                    continue  # Nächstes Encoding probieren

            return "\n".join(text_parts)

        except Exception as e:
            print(f"⚠️ CSV-Fehler: {e}")
            raise

    # ===== PowerPoint =====
    def _extract_powerpoint(self, filepath: str, max_chars: Optional[int] = None) -> str:
        """PowerPoint-Präsentationen - bricht bei max_chars ab."""
        try:
            prs = Presentation(filepath)
            text_parts = []
            total_chars = 0

            for slide_num, slide in enumerate(prs.slides):
                if max_chars and total_chars >= max_chars:
                    break

                slide_header = f"--- Slide {slide_num + 1} ---"

                if max_chars:
                    remaining = max_chars - total_chars
                    if len(slide_header) > remaining:
                        text_parts.append(slide_header[:remaining])
                        total_chars = max_chars
                    else:
                        text_parts.append(slide_header)
                        total_chars += len(slide_header)
                else:
                    text_parts.append(slide_header)
                    total_chars += len(slide_header)

                for shape in slide.shapes:
                    if max_chars and total_chars >= max_chars:
                        break

                    if hasattr(shape, "text") and shape.text.strip():
                        if max_chars:
                            remaining = max_chars - total_chars
                            if len(shape.text) > remaining:
                                text_parts.append(shape.text[:remaining])
                                total_chars = max_chars
                            else:
                                text_parts.append(shape.text)
                                total_chars += len(shape.text)
                        else:
                            text_parts.append(shape.text)
                            total_chars += len(shape.text)

                    # Tabellen in Slides
                    if hasattr(shape, "table"):
                        for row in shape.table.rows:
                            if max_chars and total_chars >= max_chars:
                                break

                            for cell in row.cells:
                                if max_chars and total_chars >= max_chars:
                                    break

                                if cell.text:
                                    if max_chars:
                                        remaining = max_chars - total_chars
                                        if len(cell.text) > remaining:
                                            text_parts.append(cell.text[:remaining])
                                            total_chars = max_chars
                                        else:
                                            text_parts.append(cell.text)
                                            total_chars += len(cell.text)
                                    else:
                                        text_parts.append(cell.text)
                                        total_chars += len(cell.text)

            return "\n".join(text_parts)

        except ImportError:
            print("⚠️ python-pptx nicht installiert. Bitte installieren: pip install python-pptx")
            raise
        except Exception as e:
            print(f"⚠️ PowerPoint-Fehler: {e}")
            raise

    # ===== OpenDocument / RTF =====
    def _extract_opendocument(self, filepath: str, max_chars: Optional[int] = None) -> str:
        """OpenDocument-Text (.odt) - bricht bei max_chars ab."""
        try:
            import zipfile
            with zipfile.ZipFile(filepath, 'r') as odt:
                content = odt.read('content.xml').decode('utf-8', errors='ignore')
                # Einfaches XML-Stripping
                import re
                text = re.sub(r'<[^>]+>', ' ', content)
                text = re.sub(r'\s+', ' ', text)

            if max_chars and len(text) > max_chars:
                text = text[:max_chars]
            return text.strip()

        except Exception as e:
            print(f"⚠️ ODT-Fehler: {e}")
            raise

    def _extract_rtf(self, filepath: str, max_chars: Optional[int] = None) -> str:
        """RTF-Dateien - bricht bei max_chars ab."""
        try:
            from striprtf.striprtf import rtf_to_text
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                rtf_content = f.read()
                text = rtf_to_text(rtf_content)

            if max_chars and len(text) > max_chars:
                text = text[:max_chars]
            return text

        except ImportError:
            # Fallback
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
                import re
                text = re.sub(r'\\[a-z]+\d*', ' ', content)
                text = re.sub(r'[{}]', ' ', text)
                text = re.sub(r'\s+', ' ', text)

            if max_chars and len(text) > max_chars:
                text = text[:max_chars]
            return text.strip()
        except Exception as e:
            print(f"⚠️ RTF-Fehler: {e}")
            raise


    def _extract_iwork_file(self, filepath: str, max_chars: Optional[int] = None) -> str:
        """
        Extrahiert Text aus Apple .pages oder .numbers Dateien (XML + IWA Fallback, Unicode-fähig, Müll reduziert).
        """
        try:
            with zipfile.ZipFile(filepath, "r") as z:

                # ===== Altes XML Format =====
                for xml_name in ["index.xml", "Index.zip"]:
                    if xml_name in z.namelist():
                        data = z.read(xml_name).decode("utf-8", errors="ignore")
                        text = re.sub(r"<[^>]+>", " ", data)
                        text = re.sub(r"\s+", " ", text)
                        return text[:max_chars] if max_chars else text

                # ===== Neues IWA Format =====
                text_parts = []

                for name in z.namelist():
                    if name.endswith(".iwa"):
                        raw = z.read(name)

                        strings = re.findall(rb"[ -~]+", raw)

                        for s in strings:
                            try:
                                text = s.decode("utf-8", errors="ignore").strip()

                                # Heuristik: mind. 3 Buchstaben/Zahlen, Sonderzeichen ≤ Buchstaben/Zahlen
                                letters_digits = sum(c.isalnum() for c in text)
                                total = len(text)
                                special_chars = total - letters_digits

                                if letters_digits >= 3 and special_chars <= letters_digits:
                                    text_parts.append(text)

                            except:
                                pass

                text = " ".join(text_parts)
                text = re.sub(r"\s+", " ", text)

                if not text:
                    raise ValueError("Kein Text im IWA gefunden")

                return text[:max_chars] if max_chars else text

        except Exception as e:
            print(f"⚠️ IWork-Datei-Fehler: {e}")
            raise

    def _extract_generic(self, filepath: str, max_chars: Optional[int] = None) -> str:
        """Fallback für nicht implementierte Formate."""
        print(f"ℹ️ Generische Extraktion für {filepath} - versuche als Textdatei")
        return self._extract_text_file(filepath, max_chars)


    def is_supported(self, filepath: str) -> bool:
        ext = os.path.splitext(filepath)[1].lower()
        return ext in self.supported_extensions

    def get_stats(self) -> dict:
        """Gibt Statistiken zurück."""
        return self.stats.copy()